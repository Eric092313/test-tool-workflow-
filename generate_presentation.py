from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_PRIMARY = RGBColor(102, 126, 234)  # #667eea
COLOR_SECONDARY = RGBColor(118, 75, 162)  # #764ba2
COLOR_TEXT = RGBColor(85, 85, 85)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_content_slide(prs, title, content_func):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    content_func(content_box)

# Slide 1: Title
add_title_slide(prs, "🎯 Inbound Software", "New Test Tool Request Workflow")

# Slide 2: Overview
def slide2_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Purpose:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    p = tf.add_paragraph()
    p.text = "Manage and streamline requests for new test tools through approval and registration"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT
    p.level = 0
    p.space_before = Pt(6)
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Key Stakeholders:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_before = Pt(12)
    
    stakeholders = [
        ("👤 DSV Test Manager", "Initiator"),
        ("🔍 QA Core", "Assessor & Researcher"),
        ("✅ QA Chapter", "Approver"),
        ("📋 EA Team", "Registrar")
    ]
    
    for icon_name, role in stakeholders:
        p = tf.add_paragraph()
        p.text = f"{icon_name} – {role}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.level = 1
        p.space_before = Pt(3)

add_content_slide(prs, "Workflow Overview", slide2_content)

# Slide 3: Process Flow
def slide3_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    flow = """DSV Test Manager Raises Request
    ↓
QA Core Reviews Request
    ↓
Tool Exists in Pool?
    ├─ YES → Provide Recommendation & Close
    └─ NO → Tool Research Phase
    ↓
QA Core Conducts Research
  • Evaluate capabilities
  • Check compatibility
  • Assess cost & licensing
    ↓
QA Chapter Reviews
    ├─ REJECTED → Close with Feedback
    └─ APPROVED
        ↓
EA Team Registers Tool
        ↓
Tool Added to Pool & Success"""
    
    p = tf.paragraphs[0]
    p.text = flow
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = 'Courier New'
    tf.word_wrap = True

add_content_slide(prs, "Process Flow Diagram", slide3_content)

# Slide 4: Stages
def slide4_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    stages = [
        ("Stage 1: Request Initiation", "DSV Test Manager raises formal request"),
        ("Stage 2: Assessment", "QA Core checks if tool exists in pool"),
        ("Stage 3: Research", "Tool research, evaluation, recommendation"),
        ("Stage 4: Approval", "QA Chapter reviews and approves/rejects"),
        ("Stage 5: Registration", "EA Team registers tool in system")
    ]
    
    p = tf.paragraphs[0]
    p.text = stages[0][0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(3)
    
    p = tf.add_paragraph()
    p.text = stages[0][1]
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.level = 1
    p.space_after = Pt(10)
    
    for stage, desc in stages[1:]:
        p = tf.add_paragraph()
        p.text = stage
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(6)
        p.space_after = Pt(3)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.level = 1
        p.space_after = Pt(6)

add_content_slide(prs, "Workflow Stages", slide4_content)

# Slide 5: Decision Points
def slide5_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Decision Point 1: Tool Pool Check"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    p = tf.add_paragraph()
    p.text = "Tool exists? YES → Fast track | NO → Research phase"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(6)
    p.space_after = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Decision Point 2: QA Chapter Approval"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Approved? YES → EA Registration | NO → Close with feedback"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(6)
    p.space_after = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Success Criteria:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_before = Pt(12)
    
    criteria = [
        "✓ Tool research is thorough and documented",
        "✓ QA Chapter approval before EA registration",
        "✓ Successfully registered in system"
    ]
    
    for criterion in criteria:
        p = tf.add_paragraph()
        p.text = criterion
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.level = 1

add_content_slide(prs, "Decision Points & Criteria", slide5_content)

# Slide 6: Timeline
def slide6_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Timeline Expectations"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    timeline = [
        ("Request Review", "2-3 business days"),
        ("Tool Research", "5-10 business days"),
        ("QA Chapter Review", "3-5 business days"),
        ("EA Registration", "2-3 business days"),
        ("TOTAL", "12-21 business days")
    ]
    
    for stage, duration in timeline:
        p = tf.add_paragraph()
        p.text = f"{stage}: {duration}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        if stage == "TOTAL":
            p.font.bold = True
            p.font.color.rgb = COLOR_SECONDARY
            p.space_before = Pt(10)
        else:
            p.level = 1
            p.space_before = Pt(4)

add_content_slide(prs, "Timeline & Expectations", slide6_content)

# Slide 7: Responsible Parties
def slide7_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Responsible Parties & Handoffs"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(12)
    
    parties = [
        ("Initiation", "DSV Test Manager", "QA Core"),
        ("Assessment", "QA Core", "QA Core (if research)"),
        ("Research", "QA Core", "QA Chapter"),
        ("Approval", "QA Chapter", "QA Core"),
        ("Registration", "QA Core + EA Team", "Complete")
    ]
    
    for stage, responsible, handoff in parties:
        p = tf.add_paragraph()
        p.text = f"{stage}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        
        p = tf.add_paragraph()
        p.text = f"Responsible: {responsible}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Hands off to: {handoff}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT
        p.level = 1
        p.space_after = Pt(6)

add_content_slide(prs, "Responsible Parties", slide7_content)

# Slide 8: Contact & Support
def slide8_content(textbox):
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Questions or Issues?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(16)
    
    contacts = [
        ("For Request Status:", "Contact QA Core"),
        ("For Tool Integration:", "Contact EA Team"),
        ("For Workflow Questions:", "Contact QA Chapter Lead")
    ]
    
    for question, contact in contacts:
        p = tf.add_paragraph()
        p.text = question
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = contact
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.level = 1

add_content_slide(prs, "Contact & Support", slide8_content)

# Save presentation
prs.save('/tmp/test-tool-workflow.pptx')
print("PowerPoint created successfully!")
